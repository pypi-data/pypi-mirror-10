#!/usr/bin/python
# -*- encoding: UTF-8 -*-
'''Survey PDF making script
'''

# built-in modules
import os
import os.path
import re
from datetime import datetime
import struct
from StringIO import StringIO
from pprint import pprint

# third-party modules
#import numpy as np
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image
import pyMTRX
import pyMTRX.experiment as om

import matplotlib.colors as mplcolors
import matplotlib.cm as cm
# Brown-Yellow color map
BrYl_cdict = {
    'red':    ((0.0,       0.0,      0.0),
               (0.0450113, 0.160784, 0.160784),
               (0.300075,  0.384926, 0.384926),
               (0.699925,  0.736303, 0.736303),
               (1.0,       1.0,      1.0)),
    'green':  ((0.0,       0.0,      0.0),
               (0.699925,  0.571275, 0.571275),
               (1.0,       1.0,      1.0)),
    'blue':   ((0.0,       0.0,      0.0),
               (0.699925,  0.0,      0.0),
               (1.0,       1.0,      1.0))
    }
BrYl = mplcolors.LinearSegmentedColormap('BrYl', BrYl_cdict)
mapper = cm.ScalarMappable(cmap=BrYl)

#==============================================================================
def main( cwd='.', sdir=None, dir_filter=range(4),
          fext='[^.()]+_mtrx', r=True, debug=False
        ):
    experiment_files = find_files(cwd, fext='mtrx', r=r)
    experiment_files.sort(key=lambda fp: os.path.getmtime(fp))
    print 'found {} .mtrx files'.format(len(experiment_files))
    for ex_fp in experiment_files:
        if not isinstance(sdir, basestring): sdir = os.path.dirname(ex_fp)
        make_pptx(
            ex_fp, sdir=sdir, dir_filter=dir_filter, fext=fext, debug=debug
        )
    # END for
    
    print 'finished'
# END main

#==============================================================================
def make_pptx( ex_fp, sdir='./', dir_filter=range(4),
               fext='[^.()]+_mtrx', debug=False
             ):
    cwd, ex_fn = os.path.split(ex_fp)
    if not cwd: cwd = '.'
    print 'working on "{}"'.format( os.path.basename(ex_fp) )
    scan_files = find_files(cwd, fext=fext)
    scan_files.sort(key=lambda fp: os.path.getmtime(fp))
    print 'found {} scan files'.format(len(scan_files))
    
    ex = om.Experiment(ex_fp, debug=debug)
    
    prs = Presentation()
    # Make title slide (starting from the blank template
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox( Pt(16), Pt(16),
                                      Inches(10)-2*Pt(16),
                                      Inches(7.5)-2*Pt(16)
                                    )
    s = ( '{1.inst_info}\n' +
          'File: {0}\n' +
          'Data Set: {1.data_set}\n' +
          'Sample: {1.sample}\n'
          'Comment: {1.comment}'
        )
    s = s.format( os.path.basename(ex_fp), ex )
    s = re.sub(r':? >', ': ', s)
    s = re.sub(r'< ', '\n', s)
    s = re.sub(r'<(?=\))', '', s)
    txBox.text_frame.text = s
    txBox.text_frame.paragraphs[0].font.size = Pt(12)
    txBox.text_frame.word_wrap = True
    
    for scn_fp in scan_files:
        if scn_fp not in ex: continue
        print 'adding {}'.format( os.path.basename(scn_fp) )
        scans = flatten( ex.import_scan(scn_fp) )
        try:
            scans = [scans[i] for i in dir_filter]
        except IndexError as err:
            print len(scans)
            print dir_filter
            raise err
        # END try
        add_slide(prs, *scans[0:2])
        add_slide(prs, *scans[2:4])
    # END for
    
    # create save name and avoid overwriting
    save_name = '{}.pptx'.format(
        re.search(r'\d{4}(?:\w{3}|\d{2})\d{2}-\d{6}', ex_fn).group(0)
    )
    i = 0
    while os.path.exists( os.path.join(cwd, save_name) ):
        i += 1
        if 99 < i: break
        save_name = re.sub(
            r'(?: \(\d\d\))?.pptx$', ' ({:02d}).pptx'.format(i), save_name
        )
    # END while
    
    prs.save(os.path.join(sdir, save_name))
    print ''
# END make__pptx

#==============================================================================
def add_slide(prs, *scans):
    if not scans: return
    # blank slide layout
    slide = prs.slides.add_slide(prs.slide_layouts[6])
        
    # Standard lengths
    # slide width and height
    #w_sld = Inches(10)
    h_sld = Inches(7.5)
    # spacer
    spc = Pt(16)
    # img/textbox width
    w_img = Pt(336)
    # textbox height
    h_txb = h_sld - (w_img + 3*spc)
    # left indents
    lindt = [spc, w_img+2*spc]
    
    # Direction translation dictionary
    scndir = { 0: 'Trace Up',   1:'Retrace Up',
               2: 'Trace Down', 3:'Retrace Down'
             }
    
    for i, scn in enumerate(scans[:2]):
        scn.global_level()
        mapper.set_array(scn.Z)
        mapper.autoscale()
        data = mapper.to_rgba(scn.Z, bytes=True)
        #data = scn.Z - np.min(scn.Z)
        #data = 255 * data / np.max(data)
        data = data.astype('uint32')
        data = data.flatten()
        bytes = struct.pack( ' '.join(len(data)*['B']), *data )
        img_buff = StringIO()
        (Image.frombytes('RGBA', scn.Z.shape, bytes)).save(
            img_buff, format='png'
        )
        
        slide.shapes.add_picture(img_buff, lindt[i], spc, width=w_img)
        txBox = slide.shapes.add_textbox(
            lindt[i], spc+w_img+spc, w_img, h_txb
        )
        tf = txBox.text_frame
        tf.text = scn.props['file']
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.bold = True
        p = tf.add_paragraph()
        p.text = '{}'.format( scndir[scn.props['direction']] )
        p.font.size = Pt(10)
        try:
            s = '{:5.3f} V\n{:5.3f} nA'.format(
                scn.props['GapVoltageControl_Voltage'].value,
                1e9 * scn.props['Regulator_Setpoint_1'].value
            )
            p = tf.add_paragraph()
            p.text = s
            p.font.size = Pt(10)
        except KeyError:
            pass
        # END try
        try:
            s = u'x= {:.1f} nm, y= {:.1f} nm, θ= {:.1f}°'.format(
                1e9 * scn.props['XYScanner_X_Offset'].value,
                1e9 * scn.props['XYScanner_Y_Offset'].value,
                scn.props['XYScanner_Angle'].value
            )
            p = tf.add_paragraph()
            p.text = s
            p.font.size = Pt(10)
        except KeyError:
            pass
        # END try
        try:
            s = u'{:.1f} × {:.1f} nm²\n{} × {} px²'.format(
                1e9 * scn.props['XYScanner_Width'].value,
                1e9 * scn.props['XYScanner_Height'].value,
                scn.props['XYScanner_Points'].value,
                scn.props['XYScanner_Lines'].value
            )
            p = tf.add_paragraph()
            p.text = s
            p.font.size = Pt(10)
        except KeyError:
            pass
        # END try
        try:
            scn.props['comment']
            tf.word_wrap = True
            p = tf.add_paragraph()
            lines = [scn.props['comment'],]
            if pyMTRX.size_change(scn):
                for t, pname, x in scn.props['pmods']:
                    lines.append(
                        '\t+{:%M:%S} | {} <-- {} {}'.format(
                            datetime.fromtimestamp(t), pname, x.value, x.unit
                        )
                    )
                # END for
            # END if
            p.text = '\n'.join(lines)
            p.font.size = Pt(10)
            p.font.italic = True
        except KeyError:
            pass
        # END try
        
        points = set()
        for crv in scn.spectra:
            x = 16 + int( 336 * crv.props['coord_px'][0] /
                          float(scn.Z.shape[0])
                        )
            y = 336+16 - int( 336 * crv.props['coord_px'][1] / 
                           float(scn.Z.shape[1])
                         )
            if (x,y) not in points:
                circ = slide.shapes.add_shape( MSO_SHAPE.OVAL,
                                            Pt(x), Pt(y), Pt(6), Pt(6)
                                          )
                circ.fill.solid()
                circ.fill.fore_color.rgb = RGBColor(0, 176, 240)
                circ.line.fill.background()
                circ.text = (
                    '{index}-{rep} {channel}\n'.format(**crv.props) #+
                    #'{:+0.3f} -> {:+0.3f} V\n'.format(crv.X[0], crv.X[-1]) +
                    #'I= {:0.3f} nA\n'.format(crv.props['Regulator_Setpoint_1'].value)
                )
                circ.text_frame.word_wrap = False
                circ.text_frame.paragraphs[0].font.size = Pt(1)
                    
                points.add((x,y))
            # END if
        # END for
    # END for
# END add_slide

#==============================================================================
def flatten(A):
    out = []
    try:
        for x in A: out.extend( flatten(x) )
        return out
    except TypeError:
        return [A]
    # END try
# END flatten

#==============================================================================
def find_files(cwd='./', fext='[^.]+', r=False):
    '''Find _mtrx files (Breath-first search)
    
    Args:
        cwd (str): current working directory
        fext (str): pattern used to match the file extensions
        r (bool): flag for recursive search
    Returns:
        (list) ['./file.ext', './file.ext', './file.ext', ...]
    '''
    
    if cwd[-1] != '/':
        cwd += '/'
    out_files = []
    work_queue = [cwd+fn for fn in os.listdir(cwd)]
    # BFS for I(t)_mtrx files
    while work_queue:
        fpath = work_queue.pop(0)
        if os.path.isdir(fpath) and r:
            work_queue.extend( [fpath+'/'+fn for fn in os.listdir(fpath)] )
        elif re.search(r'\.'+fext+'$', fpath):
            out_files.append(fpath)
        # END if
    # END while
    return out_files
# END find files

#==============================================================================
if __name__ == '__main__':
    main()
# END if
